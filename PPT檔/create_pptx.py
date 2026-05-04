from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG       = RGBColor(0xF5, 0xE6, 0xCC)
TITLE_BG = RGBColor(0x8B, 0x5E, 0x3C)
ACCENT   = RGBColor(0xD4, 0x8B, 0x3B)
DARK     = RGBColor(0x3E, 0x2C, 0x1C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xFA, 0xF0, 0xE0)
CARD_BG  = RGBColor(0xEE, 0xD9, 0xB8)
GREEN    = RGBColor(0x4C, 0x8C, 0x4A)
RED      = RGBColor(0xC0, 0x39, 0x2B)
BLUE     = RGBColor(0x2E, 0x6B, 0x9E)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color, border_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, left, top, w, h, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft JhengHei"
    p.alignment = align
    return txBox

def add_lines(slide, left, top, w, h, lines, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt, bld, clr = item[0], item[1] if len(item)>1 else False, item[2] if len(item)>2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(font_size * 0.3)
    return txBox

def slide_header(slide, title):
    set_bg(slide, BG)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), TITLE_BG)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8), title, 36, WHITE, True)

# ============================================================
# 1. 封面
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, TITLE_BG)
add_shape(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), LIGHT, ACCENT)
add_text(sl, Inches(2), Inches(2), Inches(9.3), Inches(1.2),
         "🏘️ 創業小鎮 VentureTown", 48, DARK, True, PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(3.3), Inches(9.3), Inches(0.8),
         "網頁策略經營遊戲  ·  遊戲設計概覽", 28, ACCENT, False, PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.4), Inches(9.3), Inches(0.8),
         "佈局設施 → 投入資源 → 轉換增值 → 達成目標", 22, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# 2. 遊戲概述
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "遊戲概述")

add_shape(sl, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(0.7), Inches(1.6), Inches(3.4), Inches(0.5), "🎯 目標", 22, ACCENT, True)
add_lines(sl, Inches(0.7), Inches(2.2), Inches(3.4), Inches(4.5),
    [
        "4×4 格子上佈局設施",
        "投入資源沿設施轉換增值",
        "累計收益 ≥ 目標即過關",
        "10 回合內未達標則失敗",
    ], 17)

add_shape(sl, Inches(4.7), Inches(1.5), Inches(3.8), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(4.9), Inches(1.6), Inches(3.4), Inches(0.5), "⏱️ 回合流程", 22, ACCENT, True)
add_lines(sl, Inches(4.9), Inches(2.2), Inches(3.4), Inches(4.5),
    [
        "每輪 10 回合",
        "每回合 2 選 1 行動",
        "放置 → 投入 → 轉換 → 收益",
        "每 3 回合觸發隨機事件",
    ], 17)

add_shape(sl, Inches(8.9), Inches(1.5), Inches(3.9), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(9.1), Inches(1.6), Inches(3.5), Inches(0.5), "💰 4 種資源", 22, ACCENT, True)
add_lines(sl, Inches(9.1), Inches(2.2), Inches(3.5), Inches(4.5),
    [
        "💰 金錢 — 基礎 & 收益形式",
        "🪨 原料 — 中間轉換材料",
        "📦 商品 — 可轉為金錢",
        "💎 鑽石 — 特殊×12 路線",
    ], 17)

# ============================================================
# 3. 核心循環
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "核心循環")

flow = [
    ("1️⃣ 行動選擇", "2 選 1 行動", BLUE),
    ("2️⃣ 放置設施", "手牌→格子", GREEN),
    ("3️⃣ 投入資源", "選方向投入", ACCENT),
    ("4️⃣ 設施轉換", "沿路增值", RED),
    ("5️⃣ 收益結算", "差值=收益", TITLE_BG),
]
for i, (t, d, c) in enumerate(flow):
    x = Inches(0.5 + i * 2.55)
    add_shape(sl, x, Inches(1.8), Inches(2.3), Inches(2.5), c)
    add_text(sl, x, Inches(2.1), Inches(2.3), Inches(0.6), t, 22, WHITE, True, PP_ALIGN.CENTER)
    add_text(sl, x + Inches(0.1), Inches(2.9), Inches(2.1), Inches(1), d, 18, WHITE, False, PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, x + Inches(2.15), Inches(2.5), Inches(0.5), Inches(0.5), "→", 30, DARK, True)

add_shape(sl, Inches(1.5), Inches(5), Inches(10.3), Inches(1.8), CARD_BG, ACCENT)
add_text(sl, Inches(1.8), Inches(5.15), Inches(9.7), Inches(0.5),
         "收益 = 離開時金錢 − 進入時金錢(1)", 24, ACCENT, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1.8), Inches(5.8), Inches(9.7), Inches(0.6),
         "資源值 1 進入 → 經設施轉換 → 增值後離開 → 差值即收益", 18, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# 4. 設施系統
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "設施系統 — 50 種設施")

# Basic chain
add_shape(sl, Inches(0.3), Inches(1.4), Inches(12.7), Inches(1.8), CARD_BG, ACCENT)
add_text(sl, Inches(0.5), Inches(1.5), Inches(3), Inches(0.4), "🔧 基礎轉換鏈", 20, ACCENT, True)
chain = ["💰金錢", "⛏️原料廠\n+1", "🔩精煉廠\n+2", "🏭工廠\n+1", "📦倉庫\n+2", "🏪商店\n+1", "💰金錢"]
colors = [BLUE, GREEN, GREEN, ACCENT, ACCENT, RED, BLUE]
for i, (txt, clr) in enumerate(zip(chain, colors)):
    x = Inches(0.5 + i * 1.75)
    add_shape(sl, x, Inches(2.0), Inches(1.55), Inches(0.95), clr)
    add_text(sl, x, Inches(2.05), Inches(1.55), Inches(0.9), txt, 13, WHITE, True, PP_ALIGN.CENTER)
    if i < 6:
        add_text(sl, x + Inches(1.45), Inches(2.15), Inches(0.4), Inches(0.4), "→", 20, DARK, True)

cats = [
    ("🎓 人力流", "8種", "人材產生/消耗/儲存\n來放大資源數值"),
    ("🔄 物流流", "7種", "路徑越長加成越高\n轉向/加速/放大"),
    ("⚓ 貿易流", "7種", "匯率/期貨/清倉\n貿易套利策略"),
    ("💣 拆遷流", "8種", "破壞設施換取\n爆發性收益"),
    ("🛍️ 商店系", "4種", "進階商店\n百貨公司佔2×2"),
    ("✨ 特殊", "4種", "嫉妒工廠(💎×12)\n方向/稅務/貿易代理"),
]
for i, (name, cnt, desc) in enumerate(cats):
    x = Inches(0.3 + i * 2.12)
    add_shape(sl, x, Inches(3.6), Inches(1.95), Inches(3.5), CARD_BG, ACCENT)
    add_text(sl, x+Inches(0.05), Inches(3.7), Inches(1.85), Inches(0.4), name, 15, ACCENT, True, PP_ALIGN.CENTER)
    add_text(sl, x+Inches(0.05), Inches(4.1), Inches(1.85), Inches(0.3), cnt, 14, DARK, True, PP_ALIGN.CENTER)
    add_text(sl, x+Inches(0.05), Inches(4.5), Inches(1.85), Inches(2), desc, 14, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# 5. 資源轉換路徑
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "資源轉換路徑設計")

add_shape(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.5), "🗺️ 路徑策略", 22, ACCENT, True)
add_lines(sl, Inches(0.7), Inches(2.3), Inches(5.4), Inches(4.5),
    [
        "資源從邊緣投入，沿直線前進",
        "設施排列順序 = 轉換路徑",
        "物流方向設施可改變走向",
        "增幅器/物流放大器疊加增值",
        "核心：讓資源多經過有效設施",
    ], 18)

add_shape(sl, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.5), "💎 鑽石路線（高風險高報酬）", 22, ACCENT, True)
add_lines(sl, Inches(7.2), Inches(2.3), Inches(5.4), Inches(4.5),
    [
        "嫉妒惡魔 → 解鎖嫉妒工廠",
        "金錢 → 鑽石（嫉妒工廠）",
        "鑽石經商店 → ×12 倍收益！",
        "風險：非鑽石進入 → 收益 -50%",
    ], 18)

# ============================================================
# 6. 合夥人系統
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "合夥人系統 — 30 位")

cats = [
    ("😈 惡魔系", "10位", "正面+負面並存\n貪婪:+50%收益/+50%目標\n傲慢:壓制惡魔負面", RED),
    ("⛏️ 基礎系", "3位", "穩定+1加成\n原料商/商店老闆/工廠主\n無負面效果", GREEN),
    ("🎓 人力流", "4位", "圍繞人材策略\n工會:3+人材→+50%\n總監:全用每個+15", BLUE),
    ("🚢 物流流", "4位", "強化路徑收益\n大亨:通過物流+3\n達人:4+格+10%", ACCENT),
    ("💹 貿易流", "3位", "匯率套利\n套利者:金→商→金+20%\n壟斷者:3商店+5%", TITLE_BG),
]
for i, (name, cnt, desc, clr) in enumerate(cats):
    x = Inches(0.3 + i * 2.56)
    add_shape(sl, x, Inches(1.4), Inches(2.36), Inches(3.5), clr)
    add_text(sl, x+Inches(0.1), Inches(1.5), Inches(2.16), Inches(0.4), f"{name} ({cnt})", 16, WHITE, True, PP_ALIGN.CENTER)
    add_lines(sl, x+Inches(0.1), Inches(2.1), Inches(2.16), Inches(2.5), desc.split("\n"), 13, WHITE)

# Unique partners
add_shape(sl, Inches(0.3), Inches(5.2), Inches(12.7), Inches(2), CARD_BG, ACCENT)
add_text(sl, Inches(0.5), Inches(5.3), Inches(12), Inches(0.4), "🌟 獨特合夥人", 18, ACCENT, True)
uniq = [
    ("👩‍💼 譚雅", "手牌交換"), ("👩‍🔧 蕾雅", "設施升級+2%"), ("🚛 阿北", "轉向+輸出"),
    ("🎖 市長", "中央格+5%"), ("🏗 大地主", "地圖→5×5"), ("💣 設施破壞者", "消滅+50收益"),
]
for i, (n, d) in enumerate(uniq):
    x = Inches(0.5 + i * 2.1)
    add_shape(sl, x, Inches(5.8), Inches(1.9), Inches(1.2), LIGHT, ACCENT)
    add_text(sl, x+Inches(0.05), Inches(5.85), Inches(1.8), Inches(0.4), n, 14, DARK, True, PP_ALIGN.CENTER)
    add_text(sl, x+Inches(0.05), Inches(6.25), Inches(1.8), Inches(0.4), d, 13, ACCENT, False, PP_ALIGN.CENTER)

# ============================================================
# 7. 事件系統
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "事件系統 — 19 種隨機事件")

add_shape(sl, Inches(0.3), Inches(1.4), Inches(12.7), Inches(1.4), CARD_BG, ACCENT)
add_lines(sl, Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
    [
        ("每 3 回合觸發 1 次  ·  觸發前藍色預告  ·  加權隨機抽選  ·  莫菲定律懲罰固定策略", True, ACCENT),
    ], 18)

evt_cats = [
    ("🎁 正面", ["設施補給：3選1", "原料出口熱：×2", "商品熱銷：+2", "譚雅禮物：+10", "蕾雅禮物：永久+5", "就業輔助：+20%目標"], GREEN),
    ("⚠️ 負面", ["原料大降：÷2", "颱風：方向限制", "叛亂：消滅角落設施", "危險廢棄物：放炸彈", "運輸異常：物流失效", "勞工保險：-人材×2%"], RED),
    ("🔀 中性", ["地震：設施滑動", "行業熱潮：±10%行", "區域效應：±10%列", "地塊共鳴：±10%區域", "莫菲定律：全打亂", "平靜的一天：無事"], BLUE),
]
for i, (title, items, color) in enumerate(evt_cats):
    x = Inches(0.3 + i * 4.25)
    add_shape(sl, x, Inches(3.1), Inches(4.05), Inches(4.1), color)
    add_text(sl, x+Inches(0.15), Inches(3.2), Inches(3.75), Inches(0.4), title, 18, WHITE, True)
    add_lines(sl, x+Inches(0.15), Inches(3.7), Inches(3.75), Inches(3.3), items, 15, WHITE)

# ============================================================
# 8. 動態難度 & 行動選項
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "動態難度 & 行動選項")

add_shape(sl, Inches(0.3), Inches(1.4), Inches(6.2), Inches(5.8), CARD_BG, ACCENT)
add_text(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5), "📊 動態難度", 22, ACCENT, True)
add_lines(sl, Inches(0.5), Inches(2.2), Inches(5.8), Inches(5),
    [
        ("自動適應玩家實力", True, DARK),
        "倍率範圍：0.6 ~ 2.5",
        "",
        ("⬆️ 上升條件", True, RED),
        "≤2 回合達標 → +25%",
        "連續瞬殺 → 額外 +10~20%",
        "收益 ≥3 倍目標 → +15%",
        "",
        ("⬇️ 下降條件", True, GREEN),
        ">8 回合達標 → -8%",
        "失敗重開 → 完全重置",
    ], 16)

add_shape(sl, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.8), CARD_BG, ACCENT)
add_text(sl, Inches(7), Inches(1.5), Inches(5.8), Inches(0.5), "🎬 行動選項（每回合 2 選 1）", 22, ACCENT, True)

actions = [
    ("🏗 購買設施", "4", "3 選 1 加入手牌"),
    ("⚡ 觸發事件", "2", "立即觸發隨機事件"),
    ("🔧 重新排列", "2", "自由移動所有設施"),
    ("🤝 招募合夥人", "5", "3 選 1 招募"),
    ("⬆ 提升數值", "3", "設施永久 +1"),
    ("✨ 輸出加成", "3", "下次特定類型 +3"),
    ("🔄 更換事件", "1", "重抽預告事件"),
    ("😴 跳過", "0", "不行動"),
]
for i, (name, cost, desc) in enumerate(actions):
    y = Inches(2.2 + i * 0.62)
    add_shape(sl, Inches(7), y, Inches(5.8), Inches(0.55), LIGHT, ACCENT)
    add_text(sl, Inches(7.1), y+Inches(0.05), Inches(2.2), Inches(0.4), name, 14, DARK, True)
    add_text(sl, Inches(9.3), y+Inches(0.05), Inches(0.7), Inches(0.4), cost, 13, RED, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(10), y+Inches(0.05), Inches(3.5), Inches(0.4), desc, 14, DARK)

# ============================================================
# 9. UI 介面
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "UI 介面設計")

# Left
add_shape(sl, Inches(0.5), Inches(1.5), Inches(2.5), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(0.6), Inches(1.6), Inches(2.3), Inches(0.4), "📋 左欄", 18, ACCENT, True, PP_ALIGN.CENTER)
add_lines(sl, Inches(0.6), Inches(2.2), Inches(2.3), Inches(4),
    ["🃏 手牌（扇形）", "👥 人材面板", "🤝 合夥人", "💡 提示"], 15)

# Center grid
add_shape(sl, Inches(3.3), Inches(1.5), Inches(6.5), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(3.4), Inches(1.6), Inches(6.3), Inches(0.4), "🏘️ 4×4 格子小鎮", 18, ACCENT, True, PP_ALIGN.CENTER)
grid_emoji = [
    ["⛏️","🔩","🏭","📦"],
    ["","⚡","","🏪"],
    ["🎓","","⚓",""],
    ["","🔄","","💣"],
]
for r in range(4):
    for c in range(4):
        x, y = Inches(4.3+c*1.1), Inches(2.5+r*1.05)
        add_shape(sl, x, y, Inches(1), Inches(0.95), LIGHT, ACCENT)
        if grid_emoji[r][c]:
            add_text(sl, x, y+Inches(0.05), Inches(1), Inches(0.85), grid_emoji[r][c], 26, DARK, False, PP_ALIGN.CENTER)
add_text(sl, Inches(5.7), Inches(2.1), Inches(1), Inches(0.4), "⬇️", 18, BLUE, False, PP_ALIGN.CENTER)
add_text(sl, Inches(3.8), Inches(3.8), Inches(0.5), Inches(0.4), "➡️", 18, BLUE, False, PP_ALIGN.CENTER)

# Right
add_shape(sl, Inches(10.1), Inches(1.5), Inches(2.8), Inches(5.5), CARD_BG, ACCENT)
add_text(sl, Inches(10.2), Inches(1.6), Inches(2.6), Inches(0.4), "📊 右欄", 18, ACCENT, True, PP_ALIGN.CENTER)
add_lines(sl, Inches(10.2), Inches(2.2), Inches(2.6), Inches(4),
    ["💰 本次收益", "📜 回合記錄", "🔚 結束回合", "👤 角色立繪"], 15)

# ============================================================
# 10. 設計特色
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(sl, "設計特色")

pillars = [
    ("🔄 資源流動", "金錢→原料→商品→金錢\n設計最佳轉換路徑", BLUE),
    ("⚖️ 風險報酬", "惡魔雙面刃 + 鑽石×12\n高風險高回報選擇", RED),
    ("🎲 隨機應變", "19 事件 + 莫菲定律\n動態難度每局不同", ACCENT),
    ("🧩 多元流派", "人力/物流/貿易/拆遷\n4 大流派多樣玩法", GREEN),
    ("📈 成長深度", "目標遞增 + 設施升級\n合夥人組合策略", TITLE_BG),
]
for i, (title, desc, clr) in enumerate(pillars):
    x = Inches(0.3 + i * 2.56)
    add_shape(sl, x, Inches(1.8), Inches(2.36), Inches(4.5), clr)
    add_text(sl, x+Inches(0.1), Inches(2.0), Inches(2.16), Inches(0.5), title, 22, WHITE, True, PP_ALIGN.CENTER)
    add_lines(sl, x+Inches(0.1), Inches(2.8), Inches(2.16), Inches(3), desc.split("\n"), 17, WHITE)

# Tech summary
add_shape(sl, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.7), CARD_BG, ACCENT)
add_text(sl, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
         "純 Vanilla JS 單檔架構  ·  暖色牛皮紙風格  ·  拖曳+扇形手牌  ·  114 條角色台詞  ·  自動存檔",
         16, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# 11. Thank You
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, TITLE_BG)
add_shape(sl, Inches(2), Inches(2.2), Inches(9.3), Inches(3), LIGHT, ACCENT)
add_text(sl, Inches(2.5), Inches(2.5), Inches(8.3), Inches(1),
         "🏘️ 創業小鎮 VentureTown", 40, DARK, True, PP_ALIGN.CENTER)
add_text(sl, Inches(2.5), Inches(3.5), Inches(8.3), Inches(0.7),
         "Thank You", 32, ACCENT, False, PP_ALIGN.CENTER)
add_text(sl, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.5),
         "50 設施  ·  30 合夥人  ·  19 事件  ·  4 流派", 20, DARK, False, PP_ALIGN.CENTER)

# Save
out = r"C:\Users\neionchien\Downloads\VT\VentureTown_GameDesign.pptx"
prs.save(out)
print(f"Saved to {out}")
